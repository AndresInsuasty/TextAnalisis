import os
import base64
from pptx import Presentation      


def make_presentation(titles,sumaries):
    # Creating presentation object 
    L = len(titles)
    root = Presentation() 

    # Creating slide layout 
    first_slide_layout = root.slide_layouts[1]  

    """ Ref for slide types:  
    0 ->  title and subtitle 
    1 ->  title and content 
    2 ->  section header 
    3 ->  two content 
    4 ->  Comparison 
    5 ->  Title only  
    6 ->  Blank 
    7 ->  Content with caption 
    8 ->  Pic with caption 
    """
    for i in range(L):
        slide = root.slides.add_slide(first_slide_layout) 
        slide.shapes.title.text = titles[i]
        slide.placeholders[1].text = sumaries[i]
    # Saving file 
    root.save("Output.pptx") 



def get_binary_file_downloader_html(bin_file, file_label='File'):
    with open(bin_file, 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{os.path.basename(bin_file)}">Download {file_label}</a>'
    return href

